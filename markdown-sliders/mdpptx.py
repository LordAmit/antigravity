#!/usr/bin/env python3
"""
mdpptx.py — Render the mdslides deck model to a native, editable PowerPoint
(.pptx) file via python-pptx. This is a second emitter that reuses the parsing
layer in mdslides.py (manifest, per-slide frontmatter, `###` regions) and maps
each layout to real PowerPoint shapes: text boxes with formatted runs, bullet
lists, tables, and embedded pictures.

Fidelity is intentionally "good enough, editable" rather than pixel-identical to
the Beamer PDF. Things that degrade: LaTeX code highlighting (rendered as plain
monospace), precise typographic quotes/kerning, and vector rules (drawn as thin
rectangles). Emoji render via PowerPoint's own font (no twemoji needed). Images
must be raster/EMF — a referenced .pdf/.svg is swapped for a .png sibling if one
exists, else skipped with a warning.

Invoked from mdslides.py via `--format pptx`; not a standalone CLI.
"""

import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# 16:9 canvas.
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.9
CONTENT_W = SLIDE_W - 2 * MARGIN
TITLE_TOP = 0.45
TITLE_H = 1.0
BODY_TOP = 1.75
BODY_H = SLIDE_H - BODY_TOP - 0.5


# ---------------------------------------------------------------------------
# Inline markdown -> runs
# ---------------------------------------------------------------------------

def parse_runs(text):
    """Parse inline markdown into a list of run dicts:
    {text, bold, italic, strike, underline, code, link}. Non-nested spans only
    (nested formatting degrades to the outer span); links carry a URL."""
    links = []

    def grab_link(m):
        links.append(m.group(2))
        return "\x00%d\x01%s\x02" % (len(links) - 1, m.group(1))
    text = re.sub(r"\[([^\]]+)\]\(([^)\s]+)\)", grab_link, text)

    # Wrap each span type in unique control-char delimiters, then walk.
    text = re.sub(r"\*\*(.+?)\*\*", "\x03\\1\x04", text)
    text = re.sub(r"__(.+?)__", "\x03\\1\x04", text)
    text = re.sub(r"~~(.+?)~~", "\x05\\1\x06", text)
    text = re.sub(r"<u>(.+?)</u>", "\x0e\\1\x0f", text, flags=re.S)
    text = re.sub(r"`([^`]+)`", "\x10\\1\x11", text)
    text = re.sub(r"(?<!\*)\*(.+?)\*(?!\*)", "\x12\\1\x13", text)
    text = re.sub(r"(?<!_)_(.+?)_(?!_)", "\x12\\1\x13", text)

    on = {"\x03": "bold", "\x05": "strike", "\x0e": "underline",
          "\x10": "code", "\x12": "italic"}
    off = {"\x04": "bold", "\x06": "strike", "\x0f": "underline",
           "\x11": "code", "\x13": "italic"}

    runs = []
    cur = {"bold": False, "italic": False, "strike": False,
           "underline": False, "code": False, "link": None}
    buf = []

    def flush():
        if buf:
            runs.append(dict(cur, text="".join(buf)))
            buf.clear()

    i, n = 0, len(text)
    while i < n:
        ch = text[i]
        if ch in on:
            flush()
            cur[on[ch]] = True
        elif ch in off:
            flush()
            cur[off[ch]] = False
        elif ch == "\x00":
            j = text.index("\x01", i)
            k = text.index("\x02", j)
            idx = int(text[i + 1:j])
            label = text[j + 1:k]
            flush()
            runs.append(dict(cur, text=label, link=links[idx]))
            i = k + 1
            continue
        else:
            buf.append(ch)
        i += 1
    flush()
    # Drop any stray control chars that slipped through.
    for r in runs:
        r["text"] = re.sub(r"[\x00-\x1f]", "", r["text"])
    return [r for r in runs if r["text"]]


# ---------------------------------------------------------------------------
# Block markdown -> a list of typed blocks
# ---------------------------------------------------------------------------

def parse_blocks(lines):
    """Group raw lines into typed blocks for the flow renderer."""
    blocks = []
    i, n = 0, len(lines)
    while i < n:
        s = lines[i].strip()
        if not s:
            i += 1
            continue
        m = re.match(r"^```(\w*)\s*$", s)
        if m:
            code = []
            i += 1
            while i < n and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            i += 1
            blocks.append({"type": "code", "lines": code})
            continue
        if s.startswith("|") and i + 1 < n and re.match(r"^\s*\|?[\s:|-]+\|?\s*$", lines[i + 1]):
            tbl = [lines[i], lines[i + 1]]
            i += 2
            while i < n and lines[i].strip().startswith("|"):
                tbl.append(lines[i])
                i += 1
            blocks.append({"type": "table", "rows": tbl})
            continue
        if s.startswith(">"):
            q = []
            while i < n and lines[i].strip().startswith(">"):
                q.append(re.sub(r"^\s*>\s?", "", lines[i]))
                i += 1
            blocks.append({"type": "quote", "text": " ".join(x.strip() for x in q if x.strip())})
            continue
        m = re.match(r"^#{4,6}\s+(.*)$", s)
        if m:
            blocks.append({"type": "subhead", "text": m.group(1)})
            i += 1
            continue
        m = re.match(r"^([-*+])\s+(.*)$", s)
        if m:
            items = []
            while i < n and re.match(r"^\s*[-*+]\s+", lines[i]):
                indent = len(lines[i]) - len(lines[i].lstrip())
                items.append((min(indent // 2, 2), re.sub(r"^\s*[-*+]\s+", "", lines[i])))
                i += 1
            blocks.append({"type": "bullets", "items": items})
            continue
        m = re.match(r"^\d+\.\s+", s)
        if m:
            items = []
            while i < n and re.match(r"^\s*\d+\.\s+", lines[i]):
                items.append((0, re.sub(r"^\s*\d+\.\s+", "", lines[i].strip())))
                i += 1
            blocks.append({"type": "ordered", "items": items})
            continue
        para = []
        while i < n and lines[i].strip() and not _is_block_start(lines[i]):
            para.append(lines[i].strip())
            i += 1
        blocks.append({"type": "para", "text": " ".join(para)})
    return blocks


def _is_block_start(line):
    s = line.strip()
    return bool(
        re.match(r"^```", s) or re.match(r"^[-*+]\s+", s) or re.match(r"^\d+\.\s+", s)
        or s.startswith(">") or re.match(r"^#{4,6}\s+", s) or s.startswith("|"))


# ---------------------------------------------------------------------------
# Theme -> colors / fonts
# ---------------------------------------------------------------------------

def _rgb(hexstr):
    return RGBColor.from_string(hexstr.lstrip("#"))


def theme_ctx(theme):
    pal = theme["palette"]
    dark = theme.get("dark", False)
    serif = theme.get("typography", {}).get("heading_tex", "serif") == "serif"
    return {
        "primary": _rgb(pal["primary"]),
        "secondary": _rgb(pal["secondary"]),
        "accent": _rgb(pal["accent"]),
        "bg": _rgb(pal["content_bg"]),
        "fg": _rgb(pal["content_fg"]),
        "muted": _rgb(pal["muted"]),
        "code_bg": _rgb(pal["code_bg"]),
        "heading": _rgb(pal["accent"] if dark else pal["primary"]),
        "darkfg": _rgb(pal["content_fg"]) if dark else RGBColor(0xFF, 0xFF, 0xFF),
        "font": "Georgia" if serif else "Arial",
        "mono": "Consolas",
    }


# ---------------------------------------------------------------------------
# Low-level shape helpers
# ---------------------------------------------------------------------------

def _set_strike(run):
    run._r.get_or_add_rPr().set("strike", "sngStrike")


def _apply_runs(paragraph, runs, colors, size_pt, base_color, font=None):
    for r in runs:
        run = paragraph.add_run()
        run.text = r["text"]
        f = run.font
        f.size = Pt(size_pt)
        f.name = colors["mono"] if r["code"] else (font or colors["font"])
        f.bold = r["bold"]
        f.italic = r["italic"]
        f.underline = r["underline"]
        f.color.rgb = colors["muted"] if r["code"] else base_color
        if r["strike"]:
            _set_strike(run)
        if r["link"]:
            run.hyperlink.address = r["link"]


def _textbox(slide, left, top, width, height, fill=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.fill.background()
    return box, tf


def _chars_per_line(width_in, pt):
    return max(8, int(width_in * 115 / pt))


def _est_height(block, width_in, pt):
    """Rough inches a block will occupy (for top-down flow placement)."""
    line_h = pt * 1.35 / 72.0
    cpl = _chars_per_line(width_in, pt)
    if block["type"] in ("para", "quote", "subhead"):
        lines = max(1, -(-len(block.get("text", "")) // cpl))
        pad = 0.28 if block["type"] == "quote" else 0.08
        return lines * line_h + pad
    if block["type"] in ("bullets", "ordered"):
        lines = sum(max(1, -(-len(t) // cpl)) for _, t in block["items"])
        return lines * line_h + 0.1
    if block["type"] == "code":
        return len(block["lines"]) * (pt * 1.3 / 72.0) + 0.3
    if block["type"] == "table":
        return 0.42 * len(block["rows"][:1] + block["rows"][2:])
    return line_h


# ---------------------------------------------------------------------------
# Block rendering (flow: stack shapes top-to-bottom in an area)
# ---------------------------------------------------------------------------

def render_flow(slide, blocks, area, colors, base_pt=18):
    left, top, width, height = area
    y = top
    for b in blocks:
        h = _est_height(b, width, base_pt)
        if b["type"] == "table":
            _add_table(slide, b["rows"], left, y, width, colors, base_pt)
        elif b["type"] == "quote":
            box, tf = _textbox(slide, left, y, width, h, fill=colors["code_bg"])
            p = tf.paragraphs[0]
            _apply_runs(p, [dict(r, italic=True) for r in parse_runs(b["text"])],
                        colors, base_pt, colors["muted"])
        elif b["type"] == "code":
            box, tf = _textbox(slide, left, y, width, h, fill=colors["code_bg"])
            for idx, cl in enumerate(b["lines"]):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = cl
                run.font.name = colors["mono"]
                run.font.size = Pt(base_pt - 3)
                run.font.color.rgb = colors["fg"]
        elif b["type"] == "subhead":
            box, tf = _textbox(slide, left, y, width, h)
            p = tf.paragraphs[0]
            _apply_runs(p, parse_runs(b["text"]), colors, base_pt + 2, colors["accent"])
            for run in p.runs:
                run.font.bold = True
        elif b["type"] in ("bullets", "ordered"):
            box, tf = _textbox(slide, left, y, width, h)
            for idx, (lvl, item) in enumerate(b["items"]):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.level = lvl
                marker = ("%d. " % (idx + 1)) if b["type"] == "ordered" else ("  " * lvl + "•  ")
                lead = p.add_run()
                lead.text = marker
                lead.font.size = Pt(base_pt)
                lead.font.name = colors["font"]
                lead.font.color.rgb = colors["accent"]
                _apply_runs(p, parse_runs(item), colors, base_pt, colors["fg"])
        else:  # para
            box, tf = _textbox(slide, left, y, width, h)
            _apply_runs(tf.paragraphs[0], parse_runs(b["text"]), colors, base_pt, colors["fg"])
        y += h + 0.12


def _add_table(slide, rows, left, top, width, colors, base_pt):
    def cells(line):
        return [c.strip() for c in line.strip().strip("|").split("|")]
    header = cells(rows[0])
    body = [cells(r) for r in rows[2:]]
    ncol = len(header)
    nrow = len(body) + 1
    gtbl = slide.shapes.add_table(nrow, ncol, Inches(left), Inches(top),
                                  Inches(min(width, 0.85 * width + 1)), Inches(0.4 * nrow)).table
    for c, txt in enumerate(header):
        cell = gtbl.cell(0, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        _apply_runs(p, parse_runs(txt), colors, base_pt - 2, colors["heading"])
        for run in p.runs:
            run.font.bold = True
    for r, rowcells in enumerate(body, start=1):
        rowcells = rowcells + [""] * (ncol - len(rowcells))
        for c in range(ncol):
            cell = gtbl.cell(r, c)
            cell.text = ""
            _apply_runs(cell.text_frame.paragraphs[0], parse_runs(rowcells[c]),
                        colors, base_pt - 2, colors["fg"])


def _title_box(slide, title, colors):
    _, tf = _textbox(slide, MARGIN, TITLE_TOP, CONTENT_W, TITLE_H)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _apply_runs(p, parse_runs(title), colors, 30, colors["heading"])
    for run in p.runs:
        run.font.bold = True


# ---------------------------------------------------------------------------
# Image handling
# ---------------------------------------------------------------------------

def _resolve_image(src, slide_dir):
    """Resolve an image path relative to the slide; swap .pdf/.svg for a .png
    sibling if present (python-pptx can't embed those). Returns a path or None."""
    if not os.path.isabs(src):
        src = os.path.normpath(os.path.join(slide_dir, src))
    if src.lower().endswith((".pdf", ".svg")):
        png = os.path.splitext(src)[0] + ".png"
        if os.path.isfile(png):
            return png
        print(f"warning: cannot embed '{os.path.basename(src)}' in pptx "
              f"(no .png sibling); skipping.", file=sys.stderr)
        return None
    if not os.path.isfile(src):
        print(f"warning: image not found: {src}", file=sys.stderr)
        return None
    return src


def _fit_picture(slide, path, cx, cy, max_w, max_h, scale=1.0):
    """Add a picture, scale to fit (max_w x max_h) times `scale`, center at cx,cy."""
    pic = slide.shapes.add_picture(path, Inches(0), Inches(0))
    nw, nh = pic.width, pic.height
    f = min(Inches(max_w) / nw, Inches(max_h) / nh) * scale
    pic.width = int(nw * f)
    pic.height = int(nh * f)
    pic.left = int(Inches(cx) - pic.width / 2)
    pic.top = int(Inches(cy) - pic.height / 2)
    return pic


# ---------------------------------------------------------------------------
# Layouts
# ---------------------------------------------------------------------------

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _dark_bg(slide, colors):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colors["primary"]


def slide_default(prs, colors, title, body_lines):
    s = _blank(prs)
    _bg(s, colors)
    if title:
        _title_box(s, title, colors)
    render_flow(s, parse_blocks(body_lines), (MARGIN, BODY_TOP, CONTENT_W, BODY_H), colors)
    return s


def _bg(slide, colors):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colors["bg"]


def slide_two_column(prs, colors, title, left_lines, right_lines):
    s = _blank(prs)
    _bg(s, colors)
    if title:
        _title_box(s, title, colors)
    colw = (CONTENT_W - 0.5) / 2
    render_flow(s, parse_blocks(left_lines), (MARGIN, BODY_TOP, colw, BODY_H), colors)
    render_flow(s, parse_blocks(right_lines), (MARGIN + colw + 0.5, BODY_TOP, colw, BODY_H), colors)
    return s


def slide_big_stat(prs, colors, title, stat, caption):
    s = _blank(prs)
    _bg(s, colors)
    if title:
        _title_box(s, title, colors)
    _, tf = _textbox(s, MARGIN, 2.6, CONTENT_W, 1.6)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _apply_runs(p, parse_runs(stat), colors, 72, colors["heading"])
    for run in p.runs:
        run.font.bold = True
    _, tf2 = _textbox(s, MARGIN, 4.4, CONTENT_W, 0.9)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    _apply_runs(p2, parse_runs(caption), colors, 22, colors["muted"])
    return s


def slide_image_side(prs, colors, title, img_path, side, scale, text_lines, caption, credit):
    s = _blank(prs)
    _bg(s, colors)
    if title:
        _title_box(s, title, colors)
    colw = (CONTENT_W - 0.6) / 2
    img_x = MARGIN if side == "left" else MARGIN + colw + 0.6
    txt_x = MARGIN + colw + 0.6 if side == "left" else MARGIN
    if img_path:
        _fit_picture(s, img_path, img_x + colw / 2, BODY_TOP + 2.0, colw, 3.6, scale)
        cy = BODY_TOP + 4.0
        if caption:
            _, tf = _textbox(s, img_x, cy, colw, 0.4)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _apply_runs(p, [dict(r, italic=True) for r in parse_runs(caption)], colors, 14, colors["muted"])
            cy += 0.4
        if credit:
            _, tf = _textbox(s, img_x, cy, colw, 0.35)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _apply_runs(p, parse_runs(credit), colors, 11, colors["muted"])
    render_flow(s, parse_blocks(text_lines), (txt_x, BODY_TOP, colw, BODY_H), colors)
    return s


def slide_dark(prs, colors, kind, title, subtitle, subsubtitle, kicker, leftover, ctx):
    s = _blank(prs)
    _dark_bg(s, colors)
    sizes = {"title": 40, "section": 34, "closing": 28}
    size = sizes.get(kind, 40)
    # Vertically-centered stack of centered text boxes.
    y = 2.0
    if kicker and kind == "section":
        _, tf = _textbox(s, MARGIN, y, CONTENT_W, 0.5)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _apply_runs(p, parse_runs(kicker.upper()), colors, 13, colors["muted"])
        y += 0.5
    if title:
        _, tf = _textbox(s, MARGIN, y, CONTENT_W, size / 40.0 + 0.7)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _apply_runs(p, parse_runs(title), colors, size, colors["darkfg"])
        for run in p.runs:
            run.font.bold = True
        y += size / 46.0 + 0.7
        # Short rule under the title.
        rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(SLIDE_W / 2 - 1.3),
                                  Inches(y), Inches(2.6), Pt(1.2))
        rule.fill.solid()
        rule.fill.fore_color.rgb = colors["darkfg"]
        rule.line.fill.background()
        rule.shadow.inherit = False
        y += 0.5
    if subtitle:
        _, tf = _textbox(s, MARGIN, y, CONTENT_W, 0.7)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _apply_runs(p, parse_runs(subtitle), colors, 22, colors["secondary"])
        y += 0.7
    if subsubtitle:
        _, tf = _textbox(s, MARGIN, y, CONTENT_W, 0.5)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _apply_runs(p, [dict(r, italic=True) for r in parse_runs(subsubtitle)], colors, 16, colors["darkfg"])
        y += 0.5
    if leftover:
        render_flow(s, parse_blocks(leftover), (MARGIN, y, CONTENT_W, 2.0), colors)
    return s


def _footer_logo(slide, logo_path):
    if not logo_path or not os.path.isfile(logo_path):
        return
    pic = slide.shapes.add_picture(logo_path, Inches(0), Inches(0))
    f = Inches(0.55) / pic.height
    pic.height = int(pic.height * f)
    pic.width = int(pic.width * f)
    pic.top = int(Inches(SLIDE_H - 0.8))
    pic.left = int(Inches(SLIDE_W - 0.4) - pic.width)


# ---------------------------------------------------------------------------
# Orchestration — mirrors the layout dispatch in mdslides.py
# ---------------------------------------------------------------------------

def build_pptx(deck_cfg, slides, theme, manifest_dir, out_path, extract_dark):
    """`extract_dark(slide)` is mdslides' dark-frame field extractor, returning
    (title, subtitle, subsubtitle, kicker, leftover_lines) so the two emitters
    agree on how title/section/closing content is classified."""
    from mdslides import region  # reuse region lookup

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    colors = theme_ctx(theme)

    logo_ref = deck_cfg.get("logo") or os.path.join(
        os.path.dirname(os.path.abspath(__file__)), "deck", "usf-logo.png")
    if not os.path.isabs(logo_ref):
        logo_ref = os.path.normpath(os.path.join(manifest_dir, logo_ref))
    if logo_ref.lower().endswith(".pdf"):
        alt = os.path.splitext(logo_ref)[0] + ".png"
        logo_ref = alt if os.path.isfile(alt) else logo_ref

    for slide in slides:
        layout = slide["layout"]
        fm = slide["frontmatter"]
        title = slide["title"] or _first_title(slide)

        if layout in ("title", "section", "closing"):
            t, sub, subsub, kicker, leftover = extract_dark(slide, layout)
            slide_dark(prs, colors, layout, t, sub, subsub, kicker, leftover, slide)
        elif layout == "two-column":
            slide_two_column(prs, colors, title,
                             region(slide, "left"), region(slide, "right"))
        elif layout == "big-stat":
            stat = fm.get("stat") or " ".join(
                l.strip() for l in region(slide, "stat", "number") if l.strip())
            cap = fm.get("caption") or " ".join(
                l.strip() for l in region(slide, "caption", "label") if l.strip())
            slide_big_stat(prs, colors, title, str(stat), str(cap))
        elif layout == "image-side":
            src = fm.get("image")
            if not src:
                for l in region(slide, "image"):
                    m = re.search(r"!\[[^\]]*\]\(([^)\s]+)\)", l)
                    if m:
                        src = m.group(1)
                        break
            img_path = _resolve_image(src, slide["dir"]) if src else None
            try:
                scale = float(fm["scale"]) if fm.get("scale") is not None else 1.0
            except (TypeError, ValueError):
                scale = 1.0
            text_lines = slide["body_lines"] or region(slide, "text")
            slide_image_side(prs, colors, title, img_path,
                             (fm.get("image_side") or "left").lower(), scale,
                             text_lines, fm.get("caption"), fm.get("credit"))
        else:
            body = list(slide["body_lines"])
            for name, lines in slide["regions"].items():
                body.append("#### " + name.capitalize())
                body.extend(lines)
            slide_default(prs, colors, title, body)

        if _wants_logo(fm):
            _footer_logo(prs.slides[-1], logo_ref)

    prs.save(out_path)
    return out_path


def _wants_logo(fm):
    val = fm.get("footer_logo", fm.get("logo"))
    if isinstance(val, str):
        return val.strip().lower() in ("true", "yes", "on", "1")
    return bool(val)


def _first_title(slide):
    for l in slide["body_lines"]:
        m = re.match(r"^##\s+(.*)$", l.strip())
        if m:
            return m.group(1).strip()
    return None
