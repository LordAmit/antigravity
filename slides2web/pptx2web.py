#!/usr/bin/env python3
"""
pptx2web.py - convert a .pptx into a single-page annotated-slides site.

Each slide becomes selectable, absolutely-positioned HTML (percent
coordinates, container-query font units), images are exported to media/,
and speaker notes become click-by-click annotations that dim when inactive.

Narration convention (speaker notes, per slide): a line containing only
`---` splits the notes into click-states. Markdown links/bold/italic/`code`
render in the notes panel.

Usage:
    python pptx2web.py deck.pptx -o outdir \
        --title "My talk" --byline "Jane Doe, University" \
        --venue "CONF 2026 - City - Date" \
        --toc "3:Part one" --toc "12:Part two"
"""
import argparse
import html
import json
import os
import re

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_PT = 12700


def esc(s):
    return html.escape(s, quote=False)


class Converter:
    def __init__(self, prs, outdir):
        self.prs = prs
        self.outdir = outdir
        self.w = prs.slide_width
        self.h = prs.slide_height
        self.w_pt = self.w / EMU_PER_PT
        os.makedirs(os.path.join(outdir, "media"), exist_ok=True)

    def px(self, emu):
        return f"{100.0 * emu / self.w:.2f}%"

    def py(self, emu):
        return f"{100.0 * emu / self.h:.2f}%"

    def cqw(self, pt):
        return f"{100.0 * pt / self.w_pt:.2f}cqw"

    def run_html(self, run, base_pt):
        pt = run.font.size.pt if run.font.size else base_pt
        css = [f"font-size:{self.cqw(pt)}"]
        if run.font.bold:
            css.append("font-weight:var(--bw,600)")
        if run.font.italic:
            css.append("font-style:italic")
        if run.font.underline:
            css.append("text-decoration:underline")
        try:
            if run.font.color and run.font.color.type is not None and run.font.color.rgb:
                css.append(f"color:#{run.font.color.rgb}")
        except (AttributeError, TypeError):
            pass
        span = f'<span style="{";".join(css)}">{esc(run.text)}</span>'
        addr = getattr(run.hyperlink, "address", None)
        if addr:
            span = f'<a href="{html.escape(addr, quote=True)}">{span}</a>'
        return span

    def para_html(self, para, base_pt):
        runs = "".join(self.run_html(r, base_pt) for r in para.runs) or "&nbsp;"
        pt = next((r.font.size.pt for r in para.runs if r.font.size), base_pt)
        css = ["margin:0 0 .3em 0", f"font-size:{self.cqw(pt)}"]
        if para.alignment == PP_ALIGN.CENTER:
            css.append("text-align:center")
        elif para.alignment == PP_ALIGN.RIGHT:
            css.append("text-align:right")
        if para.level and para.level > 0:
            css.append(f"margin-left:{2.2 * para.level:.1f}cqw")
        return f'<p style="{";".join(css)}">{runs}</p>'

    def shape_pos(self, sh):
        return (f"left:{self.px(sh.left)};top:{self.py(sh.top)};"
                f"width:{self.px(sh.width)};height:{self.py(sh.height)};")

    def placeholder_class(self, sh):
        try:
            ph = sh.placeholder_format
            if ph is not None and ph.type is not None:
                t = str(ph.type)
                if "TITLE" in t:
                    return "k-title"
                if "BODY" in t or "OBJECT" in t:
                    return "k-body"
        except ValueError:
            pass
        return "k-other"

    def fill_css(self, sh):
        css = ""
        try:
            if sh.fill.type is not None and str(sh.fill.type).startswith("SOLID"):
                css += f"background:#{sh.fill.fore_color.rgb};"
        except (AttributeError, TypeError):
            pass
        try:
            line = sh.line
            if line.color and line.color.type is not None and line.color.rgb:
                w = line.width.pt if line.width else 1.0
                css += f"border:{self.cqw(w)} solid #{line.color.rgb};border-radius:0.6cqw;"
        except (AttributeError, TypeError):
            pass
        return css

    def textbox_html(self, sh, slide_no):
        kls = ["shp", "txt", self.placeholder_class(sh)]
        base_pt = 28 if kls[-1] == "k-title" else 18
        try:
            if sh.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE:
                kls.append("anc-c")
            elif sh.text_frame.vertical_anchor == MSO_ANCHOR.BOTTOM:
                kls.append("anc-b")
        except (AttributeError, ValueError):
            pass
        style = self.shape_pos(sh) + "padding:0.67cqh 0.75cqw;"
        fill = self.fill_css(sh)
        body = "".join(self.para_html(p, base_pt) for p in sh.text_frame.paragraphs)
        return (f'<div data-spid="{sh.shape_id}" class="{" ".join(kls)}" '
                f'style="{style}{fill}">{body}</div>')

    def picture_html(self, sh, slide_no):
        img = sh.image
        ext = img.ext or "png"
        fname = f"slide{slide_no:02d}-{sh.shape_id}.{ext}"
        with open(os.path.join(self.outdir, "media", fname), "wb") as f:
            f.write(img.blob)
        return (f'<img data-spid="{sh.shape_id}" class="shp" '
                f'style="{self.shape_pos(sh)}" src="media/{fname}" alt="">')

    def table_html(self, sh, slide_no):
        tbl = sh.table
        total = sum(col.width for col in tbl.columns) or 1
        cols = "".join(f'<col style="width:{100.0*col.width/total:.2f}%">'
                       for col in tbl.columns)
        rows = []
        for row in tbl.rows:
            tds = []
            for cell in row.cells:
                body = "".join(self.para_html(p, 14) for p in cell.text_frame.paragraphs)
                tds.append(f'<td style="vertical-align:middle">{body}</td>')
            rows.append("<tr>" + "".join(tds) + "</tr>")
        return (f'<div data-spid="{sh.shape_id}" class="shp" style="{self.shape_pos(sh)}">'
                f'<table>{cols}{"".join(rows)}</table></div>')

    def shape_html(self, sh, slide_no):
        st = sh.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            return self.picture_html(sh, slide_no)
        if st == MSO_SHAPE_TYPE.TABLE:
            return self.table_html(sh, slide_no)
        if st == MSO_SHAPE_TYPE.GROUP:
            return "".join(self.shape_html(s, slide_no) for s in sh.shapes)
        if sh.has_text_frame and sh.text_frame.text.strip():
            return self.textbox_html(sh, slide_no)
        if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            fill = self.fill_css(sh)
            if fill:
                return (f'<div data-spid="{sh.shape_id}" class="shp" '
                        f'style="{self.shape_pos(sh)}{fill}"></div>')
        return ""

    def slide_html(self, slide, slide_no):
        inner = "".join(self.shape_html(sh, slide_no) for sh in slide.shapes)
        return f'<div class="slide" id="s{slide_no}">{inner}</div>'


def notes_segments(slide):
    if not slide.has_notes_slide:
        return [""]
    txt = slide.notes_slide.notes_text_frame.text or ""
    segs = [s.strip() for s in re.split(r"\n\s*---\s*\n", txt)]
    return segs if segs else [""]


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("-o", "--outdir", default="site")
    ap.add_argument("--title", default=None)
    ap.add_argument("--byline", default="")
    ap.add_argument("--venue", default="")
    ap.add_argument("--toc", action="append", default=[],
                    help='e.g. --toc "6:Part 1" (slide number : label)')
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    os.makedirs(args.outdir, exist_ok=True)
    conv = Converter(prs, args.outdir)

    frags, narr = {}, {}
    for i, slide in enumerate(prs.slides, start=1):
        frags[str(i)] = conv.slide_html(slide, i)
        narr[str(i)] = notes_segments(slide)
    n = len(prs.slides)

    toc_html = ""
    if args.toc:
        entries = []
        for t in args.toc:
            num, label = t.split(":", 1)
            entries.append((int(num), label.strip()))
        entries.sort()
        bounds = [e[0] for e in entries] + [n + 1]
        toc_html = '<div class="t">In this talk</div>' + "".join(
            f'<a data-lo="{lo}" data-hi="{hi}" onclick="jump({lo})">{esc(label)}</a>'
            for (lo, label), hi in zip(entries, bounds[1:])
        )

    tpl_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                            "viewer_template.html")
    with open(tpl_path, encoding="utf-8") as f:
        tpl = f.read()

    title = args.title or os.path.splitext(os.path.basename(args.pptx))[0]
    aspect = f"{prs.slide_width} / {prs.slide_height}"
    out = (tpl.replace("{{TITLE}}", esc(title))
              .replace("{{BYLINE}}", esc(args.byline))
              .replace("{{VENUE}}", esc(args.venue))
              .replace("{{ASPECT}}", aspect)
              .replace("{{TOC}}", toc_html)
              .replace("{{FRAGS}}", json.dumps(frags, ensure_ascii=False))
              .replace("{{NARR}}", json.dumps(narr, ensure_ascii=False))
              .replace("{{N}}", str(n)))

    with open(os.path.join(args.outdir, "index.html"), "w", encoding="utf-8") as f:
        f.write(out)
    print(f"Wrote {args.outdir}/index.html ({n} slides), media/ has "
          f"{len(os.listdir(os.path.join(args.outdir, 'media')))} file(s).")


if __name__ == "__main__":
    main()
