"""Build a small demo deck exercising every converter feature:
placeholders, bullets, hyperlink, picture, rounded rect, table,
click animations (injected timing XML), and [click]-split notes."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw
from lxml import etree

# a small chart-ish PNG
img = Image.new("RGB", (640, 400), "white")
d = ImageDraw.Draw(img)
for i, h in enumerate([120, 180, 260, 330]):
    x = 60 + i * 140
    d.rectangle([x, 380 - h, x + 90, 380], fill=(245, 106, 106))
d.line([40, 380, 620, 380], fill="black", width=3)
d.line([40, 20, 40, 380], fill="black", width=3)
img.save("fig.png")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- slide 1: title ---
s = prs.slides.add_slide(prs.slide_layouts[0])
s.shapes.title.text = "Annotated slides, as real HTML"
s.placeholders[1].text = "A demo deck for slides2web"
s.notes_slide.notes_text_frame.text = (
    "Welcome. This page was generated straight from a .pptx file - the text "
    "you see in the slide above is selectable HTML, not a screenshot. "
    "Use the arrow keys, or click Next.")

# --- slide 2: bullets + picture + callout, with two click animations ---
s = prs.slides.add_slide(prs.slide_layouts[5])
s.shapes.title.text = "Click-by-click builds"
body = s.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(6.2), Inches(4.5))
tf = body.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "This bullet is visible from the start."
p.runs[0].font.size = Pt(20)
p2 = tf.add_paragraph()
p2.text = "Shapes animated in PowerPoint appear on their click."
p2.runs[0].font.size = Pt(20)
pic = s.shapes.add_picture("fig.png", Inches(7.3), Inches(1.7), Inches(5.2), Inches(3.3))
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(7.3), Inches(5.3), Inches(5.2), Inches(1.3))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xDE, 0xEB, 0xF6)
box.line.color.rgb = RGBColor(0x8E, 0xAA, 0xDB)
btf = box.text_frame
btf.text = "This callout arrives on the second click."
btf.paragraphs[0].runs[0].font.size = Pt(16)
r = btf.paragraphs[0].add_run()
r.text = " Learn more"
r.font.size = Pt(16)
r.hyperlink.address = "https://amitsealami.com"
s.notes_slide.notes_text_frame.text = (
    "A slide can build up over several clicks, exactly as delivered.\n"
    "[click]\n"
    "First click: the **chart** fades in. Its bars are just a picture "
    "exported from the deck into media/.\n"
    "[click]\n"
    "Second click: the callout box appears, with a working *hyperlink*. "
    "Each click also advances this annotation panel - see [the original]"
    "(https://www.cs.princeton.edu/~arvindn/talks/icml-2026-annotated-slides/) "
    "that inspired the format.")

# inject minimal PowerPoint timing XML: pic on click 1, box on click 2
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
def spTgt(spid):
    return ('<p:tgtEl xmlns:p="%s"><p:spTgt spid="%d"/></p:tgtEl>' % (P, spid))
timing = etree.fromstring(
    '<p:timing xmlns:p="%s"><p:tnLst><p:par>'
    '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
    '<p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
    '<p:childTnLst>'
    '<p:par><p:cTn id="3" fill="hold"><p:childTnLst><p:par><p:cTn id="4" nodeType="clickEffect">'
    '<p:childTnLst><p:set><p:cBhvr>%s</p:cBhvr></p:set></p:childTnLst>'
    '</p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
    '<p:par><p:cTn id="5" fill="hold"><p:childTnLst><p:par><p:cTn id="6" nodeType="clickEffect">'
    '<p:childTnLst><p:set><p:cBhvr>%s</p:cBhvr></p:set></p:childTnLst>'
    '</p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
    '</p:childTnLst></p:cTn></p:seq>'
    '</p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
    % (P, spTgt(pic.shape_id), spTgt(box.shape_id)))
s.element.append(timing)

# --- slide 3: table ---
s = prs.slides.add_slide(prs.slide_layouts[5])
s.shapes.title.text = "Tables survive too"
tbl = s.shapes.add_table(3, 3, Inches(1.5), Inches(2.0),
                         Inches(10.3), Inches(3.0)).table
heads = ["Stage", "PowerPoint", "On the web"]
rows = [["Text", "Text boxes", "Positioned divs, cqw-sized"],
        ["Builds", "Click animations", "data-step reveals"]]
for j, h in enumerate(heads):
    tbl.cell(0, j).text = h
for i, row in enumerate(rows, start=1):
    for j, v in enumerate(row):
        tbl.cell(i, j).text = v
s.notes_slide.notes_text_frame.text = (
    "Tables come across as real <table> elements with proportional column "
    "widths, so they reflow crisply at any window size.\n"
    "[click]\n"
    "Notes can have more segments than the slide has animations - extra "
    "segments simply become extra clicks. That is what just happened.")

prs.save("demo.pptx")
print("demo.pptx written")
